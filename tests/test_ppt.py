# tests/test_ppt.py
import sys
import os
import pytest
sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..', 'src'))

from pptx import Presentation
import theme

PPT_PATH = os.path.join(os.path.dirname(__file__), '..', 'output', 'network-card-csi-v4.0.pptx')

@pytest.fixture(scope="module")
def prs():
    assert os.path.exists(PPT_PATH), f"PPT not found: {PPT_PATH}"
    return Presentation(PPT_PATH)

def test_slide_count(prs):
    import sys, os
    sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..', 'src'))
    import content as _content
    assert len(prs.slides) == len(_content.SLIDES)

def test_slide_dimensions(prs):
    assert prs.slide_width == theme.SLIDE_WIDTH
    assert prs.slide_height == theme.SLIDE_HEIGHT

def test_table_slides_exist(prs):
    """Verify that table slides exist"""
    table_slides = [s for s in prs.slides
                    if any(sh.has_table for sh in s.shapes)]
    assert len(table_slides) >= 4, "Should have at least 4 table slides"

def test_all_slides_have_shapes(prs):
    for i, slide in enumerate(prs.slides):
        assert len(slide.shapes) > 0, f"Slide {i+1} has no shapes"

def test_no_offscreen_shapes(prs):
    """確認所有 shape 都在投影片可見範圍內"""
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            bottom = shape.top + shape.height
            right = shape.left + shape.width
            assert bottom <= theme.SLIDE_HEIGHT + 1000, (
                f"Slide {i+1} shape '{shape.name}' bottom={bottom} exceeds SLIDE_HEIGHT={theme.SLIDE_HEIGHT}"
            )
            assert right <= theme.SLIDE_WIDTH + 1000, (
                f"Slide {i+1} shape '{shape.name}' right={right} exceeds SLIDE_WIDTH={theme.SLIDE_WIDTH}"
            )
