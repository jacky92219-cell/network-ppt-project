# src/builders.py
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
import theme


# ─────────────────────── XML helpers ───────────────────────

def _add_arrowhead(connector):
    """在 connector 末端加箭頭"""
    sp = connector._element
    spPr = sp.find(qn('p:spPr'))
    if spPr is None:
        return
    ln = spPr.find(qn('a:ln'))
    if ln is None:
        return
    for old in ln.findall(qn('a:tailEnd')):
        ln.remove(old)
    tailEnd = OxmlElement('a:tailEnd')
    tailEnd.set('type', 'arrow')
    tailEnd.set('w', 'med')
    tailEnd.set('len', 'med')
    ln.append(tailEnd)


def _darken(color: RGBColor, factor: float = 0.5) -> RGBColor:
    """把 RGBColor 調暗（factor=0.5 → 50% 亮度）"""
    return RGBColor(
        int(color[0] * factor),
        int(color[1] * factor),
        int(color[2] * factor),
    )


def _set_rounded_corner(shape, val: int = 16667):
    """設定圓角矩形的圓角程度（type 5）"""
    adj = shape._element.spPr.find(qn('a:prstGeom'))
    if adj is None:
        return
    avLst = adj.find(qn('a:avLst'))
    if avLst is None:
        return
    for gd in avLst.findall(qn('a:gd')):
        avLst.remove(gd)
    gd_el = OxmlElement('a:gd')
    gd_el.set('name', 'adj')
    gd_el.set('fmla', f'val {val}')
    avLst.append(gd_el)


# ─────────────────────── 基礎元件 ───────────────────────

def set_slide_background(slide, color: RGBColor):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, text, left, top, width, height,
                font_name=theme.FONT_BODY, font_size=theme.BODY_SIZE,
                color=theme.TEXT_COLOR, bold=False, align=PP_ALIGN.LEFT,
                word_wrap=True):
    txBox = slide.shapes.add_textbox(
        int(left), int(top), int(width), int(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    run.font.color.rgb = color
    run.font.bold = bold
    return txBox


# ─────────────────────── Grid-aware 共用裝飾元件 ───────────────────────

def add_title_bar(slide, title_text: str, section: int = 0):
    """全寬標題列：暗色背景 + section 色底線 + 白色標題文字"""
    bar_h = int(theme.TITLE_BAR_H)
    bar_color = theme.TITLE_BAR_COLORS.get(section, theme.TITLE_BAR_COLORS[0])
    sec_color = theme.SECTION_COLORS.get(section, theme.PRIMARY_COLOR) if section > 0 else theme.PRIMARY_COLOR

    bar = slide.shapes.add_shape(1, 0, 0, int(theme.SLIDE_WIDTH), bar_h)
    # 水平漸層：左暗（title bar色）→ 右微亮（section color暗版）
    bar.fill.gradient()
    bar.fill.gradient_angle = 0.0  # 0° = left to right
    stops = bar.fill.gradient_stops
    stops[0].position = 0.0
    stops[0].color.rgb = bar_color
    stops[1].position = 1.0
    stops[1].color.rgb = _darken(sec_color, 0.30)
    bar.line.fill.background()

    # 底部 section-color 亮線
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, 0, bar_h, int(theme.SLIDE_WIDTH), bar_h
    )
    connector.line.color.rgb = sec_color
    connector.line.width = Pt(2.5)

    # 標題文字
    add_textbox(slide, title_text,
                int(theme.CONTENT_LEFT),
                int(Inches(0.1)),
                int(theme.CONTENT_WIDTH),
                bar_h - int(Inches(0.1)),
                font_name=theme.FONT_TITLE, font_size=theme.TITLE_SIZE,
                color=theme.TEXT_COLOR, bold=True)


def add_content_panel(slide, top=None, height=None, left=None, width=None):
    """深色圓角矩形內容面板（Glassmorphism 漸層）"""
    if left is None:
        left = int(theme.CONTENT_LEFT)
    if width is None:
        width = int(theme.CONTENT_WIDTH)
    if top is None:
        top = int(theme.CONTENT_TOP)
    if height is None:
        height = int(theme.CONTENT_HEIGHT)

    panel = slide.shapes.add_shape(5, int(left), int(top), int(width), int(height))
    # Glassmorphism: 垂直漸層（頂部較亮 → 底部較深）
    panel.fill.gradient()
    panel.fill.gradient_angle = 270.0  # 270° = top to bottom
    stops = panel.fill.gradient_stops
    stops[0].position = 0.0
    stops[0].color.rgb = RGBColor(0x14, 0x1c, 0x27)  # 頂部：比 PANEL_COLOR 亮
    stops[1].position = 1.0
    stops[1].color.rgb = theme.PANEL_COLOR              # 底部：標準深色
    panel.line.color.rgb = theme.PANEL_BORDER
    panel.line.width = Pt(0.75)
    _set_rounded_corner(panel, 10000)

    # Top highlight strip（模擬玻璃頂部反光）
    highlight = slide.shapes.add_shape(1,
        int(left) + 2, int(top) + 2,
        int(width) - 4, int(Inches(0.025)))
    highlight.fill.solid()
    highlight.fill.fore_color.rgb = RGBColor(0x3a, 0x4a, 0x5c)
    highlight.line.fill.background()

    return panel


def add_footer_bar(slide, number: int, section: int = 0):
    """底部資訊列：段落名稱 + 頁碼"""
    footer_h = int(theme.FOOTER_H)
    footer_y = int(theme.SLIDE_HEIGHT - theme.FOOTER_H)
    sec_color = theme.SECTION_COLORS.get(section, theme.SUBTEXT_COLOR) if section > 0 else theme.PRIMARY_COLOR
    sec_name = theme.SECTION_NAMES.get(section, "")

    footer_bg = slide.shapes.add_shape(1, 0, footer_y, int(theme.SLIDE_WIDTH), footer_h)
    footer_bg.fill.solid()
    footer_bg.fill.fore_color.rgb = theme.PANEL_COLOR
    footer_bg.line.fill.background()

    # 頂部 section-color 細線
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, 0, footer_y, int(theme.SLIDE_WIDTH), footer_y
    )
    connector.line.color.rgb = sec_color
    connector.line.width = Pt(1.0)

    text_h = int(footer_h - Inches(0.1))
    if sec_name:
        add_textbox(slide, sec_name,
                    int(Inches(0.3)), footer_y + int(Inches(0.05)),
                    int(Inches(5.0)), text_h,
                    font_size=Pt(11), color=sec_color)

    add_textbox(slide, str(number),
                int(theme.CONTENT_RIGHT - Inches(0.8)),
                footer_y + int(Inches(0.05)),
                int(Inches(0.8)), text_h,
                font_size=Pt(11), color=theme.SUBTEXT_COLOR,
                align=PP_ALIGN.RIGHT)


def add_note(slide, note_text):
    """備註文字（使用 grid 常數，確保不與 footer 重疊）"""
    add_textbox(slide, f"▶ {note_text}",
                int(theme.CONTENT_LEFT + Inches(0.1)),
                int(theme.NOTE_TOP),
                int(theme.CONTENT_WIDTH - Inches(0.2)),
                int(theme.NOTE_HEIGHT),
                font_size=theme.SMALL_SIZE, color=theme.SUBTEXT_COLOR)


# ─────────────────────── Builders ───────────────────────

def build_cover(slide, data):
    set_slide_background(slide, theme.BG_COLOR)
    section = data.get("section", 0)
    sw = int(theme.SLIDE_WIDTH)
    sh = int(theme.SLIDE_HEIGHT)

    # 幾何裝飾色塊（四角分佈，各用不同 section color）
    # 左上大色塊（section 1 電光藍）
    blk = slide.shapes.add_shape(1, 0, 0, int(Inches(2.5)), int(Inches(2.0)))
    blk.fill.solid()
    blk.fill.fore_color.rgb = _darken(theme.SECTION_COLORS[1], 0.45)
    blk.line.fill.background()

    # 右上三角形（type 6 = right triangle，section 2 紫）
    tri = slide.shapes.add_shape(6,
        int(sw - Inches(2.2)), 0, int(Inches(2.2)), int(Inches(2.0)))
    tri.fill.solid()
    tri.fill.fore_color.rgb = _darken(theme.SECTION_COLORS[2], 0.45)
    tri.line.fill.background()

    # 右中菱形（type 4 = diamond，section 3 青綠）
    dia = slide.shapes.add_shape(4,
        int(sw - Inches(1.5)), int(Inches(2.1)), int(Inches(1.3)), int(Inches(1.1)))
    dia.fill.solid()
    dia.fill.fore_color.rgb = _darken(theme.SECTION_COLORS[3], 0.45)
    dia.line.fill.background()

    # 右下小方塊（section 4 金）
    sq = slide.shapes.add_shape(1,
        int(sw - Inches(0.7)), int(sh - Inches(0.7)), int(Inches(0.6)), int(Inches(0.6)))
    sq.fill.solid()
    sq.fill.fore_color.rgb = _darken(theme.SECTION_COLORS[4], 0.5)
    sq.line.fill.background()

    # 全寬螢光青綠粗線（標題上方）
    line_y = int(Inches(1.55))
    top_line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, 0, line_y, sw, line_y)
    top_line.line.color.rgb = theme.ACCENT2_COLOR
    top_line.line.width = Pt(4.0)

    # 主標題
    add_textbox(slide, data["title"],
                int(Inches(0.3)), int(Inches(1.65)), int(Inches(9.4)), int(Inches(1.5)),
                font_name=theme.FONT_TITLE, font_size=Pt(48),
                color=theme.PRIMARY_COLOR, bold=True, align=PP_ALIGN.CENTER)

    # 全寬鮮紅粗線（分隔標題/副標題）
    red_y = int(Inches(3.25))
    red_line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, 0, red_y, sw, red_y)
    red_line.line.color.rgb = theme.ACCENT_COLOR
    red_line.line.width = Pt(4.0)

    # 標題左側 accent bar
    for ax in [int(Inches(0.3)), int(sw - Inches(0.35))]:
        acc = slide.shapes.add_shape(1,
            ax, int(Inches(1.7)), int(Inches(0.05)), int(Inches(1.4)))
        acc.fill.solid()
        acc.fill.fore_color.rgb = theme.PRIMARY_COLOR
        acc.line.fill.background()

    # 副標題
    add_textbox(slide, data["subtitle"],
                int(Inches(0.4)), int(Inches(3.32)), int(Inches(9.2)), int(Inches(1.0)),
                font_size=theme.SUBTITLE_SIZE,
                color=theme.ACCENT2_COLOR, align=PP_ALIGN.CENTER)

    # 日期 + 版號（y=4.45"，確保在畫面內）
    date_ver = data["date"]
    if data.get("version"):
        date_ver = f"{data['date']}　　{data['version']}"
    add_textbox(slide, date_ver,
                int(Inches(0.4)), int(Inches(4.45)), int(Inches(9.2)), int(Inches(0.45)),
                font_size=theme.SMALL_SIZE,
                color=theme.SUBTEXT_COLOR, align=PP_ALIGN.CENTER)

    if data.get("_slide_num"):
        add_footer_bar(slide, data["_slide_num"], section)


def build_section_break(slide, data):
    """段落過場：頂條 + 中央面板 + accent bar"""
    section = data.get("section", 1)
    sec_color = theme.SECTION_COLORS.get(section, theme.PRIMARY_COLOR)
    dark_color = _darken(sec_color, 0.18)
    sw = int(theme.SLIDE_WIDTH)
    sh = int(theme.SLIDE_HEIGHT)

    set_slide_background(slide, dark_color)

    # 頂部全寬 section-color 粗條
    top_bar = slide.shapes.add_shape(1, 0, 0, sw, int(Inches(0.18)))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = sec_color
    top_bar.line.fill.background()

    # 中央面板
    panel_w = int(Inches(8.4))
    panel_h = int(Inches(2.5))
    panel_x = int((theme.SLIDE_WIDTH - panel_w) // 2)
    panel_y = int(Inches(1.2))
    panel = slide.shapes.add_shape(1, panel_x, panel_y, panel_w, panel_h)
    panel.fill.solid()
    panel.fill.fore_color.rgb = _darken(sec_color, 0.28)
    panel.line.color.rgb = sec_color
    panel.line.width = Pt(2.0)

    # 左側粗 accent bar
    acc_bar = slide.shapes.add_shape(1,
        panel_x, panel_y, int(Inches(0.1)), panel_h)
    acc_bar.fill.solid()
    acc_bar.fill.fore_color.rgb = sec_color
    acc_bar.line.fill.background()

    # 段落大標題
    add_textbox(slide, data["title"],
                panel_x + int(Inches(0.2)), panel_y,
                panel_w - int(Inches(0.2)), int(Inches(1.5)),
                font_name=theme.FONT_TITLE, font_size=Pt(44),
                color=theme.TEXT_COLOR, bold=True, align=PP_ALIGN.CENTER)

    # 副標題
    add_textbox(slide, data.get("subtitle", ""),
                panel_x + int(Inches(0.2)),
                panel_y + int(Inches(1.55)),
                panel_w - int(Inches(0.2)),
                int(Inches(0.85)),
                font_size=Pt(18),
                color=sec_color, align=PP_ALIGN.CENTER)

    # 底部全寬水平線
    bottom_line_y = int(Inches(4.6))
    b_line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, 0, bottom_line_y, sw, bottom_line_y)
    b_line.line.color.rgb = sec_color
    b_line.line.width = Pt(1.5)

    # "Section N" 標籤（確保在 4.8"，畫面內）
    sec_num_text = f"Section {section}"
    add_textbox(slide, sec_num_text,
                int(Inches(0.4)), int(Inches(4.72)),
                int(Inches(2.5)), int(Inches(0.35)),
                font_size=Pt(12), color=sec_color)

    if data.get("_slide_num"):
        add_footer_bar(slide, data["_slide_num"], section)


def build_bullets(slide, data):
    set_slide_background(slide, theme.BG_COLOR)
    section = data.get("section", 0)
    add_title_bar(slide, data["title"], section)
    add_content_panel(slide)

    sec_color = theme.SECTION_COLORS.get(section, theme.PRIMARY_COLOR) if section > 0 else theme.PRIMARY_COLOR

    txBox = slide.shapes.add_textbox(
        int(theme.CONTENT_LEFT + Inches(0.15)),
        int(theme.CONTENT_TOP + Inches(0.1)),
        int(theme.CONTENT_WIDTH - Inches(0.3)),
        int(theme.CONTENT_HEIGHT - theme.NOTE_HEIGHT - Inches(0.15))
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True

    for bullet in data["bullets"]:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        if bullet == "":
            run.text = bullet
            run.font.size = Pt(6)
            p.space_before = Pt(2)
        elif bullet.startswith("  "):
            run.text = "  › " + bullet.lstrip()
            run.font.size = Pt(15)
            run.font.color.rgb = theme.SUBTEXT_COLOR
            p.level = 1
            p.space_before = Pt(2)
        else:
            if bullet.startswith("▶") or bullet.endswith("：") or (len(bullet) > 1 and bullet[0].isdigit()):
                run.text = bullet
                run.font.bold = True
                run.font.color.rgb = sec_color
            else:
                run.text = "● " + bullet
                run.font.color.rgb = theme.TEXT_COLOR
            run.font.size = theme.BODY_SIZE
            p.space_before = Pt(5)

    if "note" in data:
        add_note(slide, data["note"])
    if data.get("_slide_num"):
        add_footer_bar(slide, data["_slide_num"], section)


def build_two_col(slide, data):
    set_slide_background(slide, theme.BG_COLOR)
    section = data.get("section", 0)
    add_title_bar(slide, data["title"], section)

    sec_color = theme.SECTION_COLORS.get(section, theme.PRIMARY_COLOR) if section > 0 else theme.ACCENT_COLOR
    col_w = int((theme.CONTENT_WIDTH - theme.GUTTER) // 2)
    left_x = int(theme.CONTENT_LEFT)
    right_x = int(theme.CONTENT_LEFT + col_w + theme.GUTTER)
    ct = int(theme.CONTENT_TOP)
    ch = int(theme.CONTENT_HEIGHT)

    add_content_panel(slide, top=ct, height=ch, left=left_x, width=col_w)
    add_content_panel(slide, top=ct, height=ch, left=right_x, width=col_w)

    # 中間分隔線
    div_x = int(left_x + col_w + theme.GUTTER // 2)
    div = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, div_x, ct, div_x, ct + ch)
    div.line.color.rgb = sec_color
    div.line.width = Pt(1.0)

    def add_col(title, bullets, col_left):
        add_textbox(slide, title,
                    col_left + int(Inches(0.15)),
                    ct + int(Inches(0.1)),
                    col_w - int(Inches(0.3)),
                    int(Inches(0.45)),
                    font_size=Pt(17), color=sec_color, bold=True)
        txBox = slide.shapes.add_textbox(
            col_left + int(Inches(0.15)),
            ct + int(Inches(0.6)),
            col_w - int(Inches(0.3)),
            ch - int(Inches(0.7))
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        first = True
        for b in bullets:
            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()
            run = p.add_run()
            run.text = b
            if b == "":
                run.font.size = Pt(6)
            elif b.startswith("  "):
                run.font.size = Pt(13)
                run.font.color.rgb = theme.SUBTEXT_COLOR
                p.space_before = Pt(2)
            else:
                run.font.size = Pt(15)
                run.font.color.rgb = theme.TEXT_COLOR
                p.space_before = Pt(4)
                if b.startswith("✅") or b.startswith("❌") or b.startswith("⚠") or b.startswith("✗"):
                    run.font.bold = True

    add_col(data["left_title"], data["left_bullets"], left_x)
    add_col(data["right_title"], data["right_bullets"], right_x)

    if "note" in data:
        add_note(slide, data["note"])
    if data.get("_slide_num"):
        add_footer_bar(slide, data["_slide_num"], section)


def build_table(slide, data):
    set_slide_background(slide, theme.BG_COLOR)
    section = data.get("section", 0)
    add_title_bar(slide, data["title"], section)
    add_content_panel(slide)

    sec_color = theme.SECTION_COLORS.get(section, theme.PRIMARY_COLOR) if section > 0 else theme.PRIMARY_COLOR
    rows = len(data["rows"]) + 1
    cols = len(data["headers"])
    left = int(theme.CONTENT_LEFT + Inches(0.05))
    top = int(theme.CONTENT_TOP + Inches(0.05))
    width = int(theme.CONTENT_WIDTH - Inches(0.1))
    max_h = int(theme.CONTENT_HEIGHT - Inches(0.5))
    height = min(int(Inches(0.5) * rows + Inches(0.1)), max_h)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    col_width = int(width / cols)
    for i in range(cols):
        table.columns[i].width = col_width if i < cols - 1 else (width - col_width * (cols - 1))

    for ci, hdr in enumerate(data["headers"]):
        cell = table.cell(0, ci)
        cell.text = hdr
        cell.fill.solid()
        cell.fill.fore_color.rgb = sec_color
        cell.margin_left = Inches(0.08)
        cell.margin_right = Inches(0.08)
        cell.margin_top = Inches(0.05)
        cell.margin_bottom = Inches(0.05)
        p = cell.text_frame.paragraphs[0]
        run = p.runs[0] if p.runs else p.add_run()
        run.font.bold = True
        run.font.size = Pt(14)
        run.font.color.rgb = RGBColor(0x0d, 0x11, 0x17)  # 深色字在亮色表頭上
        run.font.name = theme.FONT_BODY
        p.alignment = PP_ALIGN.CENTER

    for ri, row in enumerate(data["rows"]):
        bg = theme.BG_COLOR if ri % 2 == 0 else theme.TABLE_ROW_ALT
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            p = cell.text_frame.paragraphs[0]
            run = p.runs[0] if p.runs else p.add_run()
            run.font.size = Pt(13)
            if ci == 0:
                run.font.bold = True
                run.font.color.rgb = sec_color
            else:
                run.font.color.rgb = theme.TEXT_COLOR
            run.font.name = theme.FONT_BODY

    if "note" in data:
        add_note(slide, data["note"])
    if data.get("_slide_num"):
        add_footer_bar(slide, data["_slide_num"], section)


def build_flow(slide, data):
    set_slide_background(slide, theme.BG_COLOR)
    section = data.get("section", 0)
    add_title_bar(slide, data["title"], section)
    add_content_panel(slide)

    sec_color = theme.SECTION_COLORS.get(section, theme.PRIMARY_COLOR) if section > 0 else theme.PRIMARY_COLOR
    dark_sec = _darken(sec_color, 0.30)

    items = data["flow_items"]
    n = len(items)
    box_w = int(Inches(1.5))
    gap = int(Inches(0.4))    # 從 0.18" 加大到 0.4"，箭頭清晰可見
    box_h = int(Inches(1.05))
    total_w = n * box_w + (n - 1) * gap
    start_x = int((theme.SLIDE_WIDTH - total_w) // 2)
    # 垂直置中在 content area
    y = int(theme.CONTENT_TOP + (theme.CONTENT_HEIGHT - box_h) // 2)

    for i, (label, desc) in enumerate(items):
        x = int(start_x + i * (box_w + gap))
        shape = slide.shapes.add_shape(5, x, y, box_w, box_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = dark_sec
        shape.line.color.rgb = sec_color
        shape.line.width = Pt(1.5)
        _set_rounded_corner(shape, 20000)

        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = int(Inches(0.06))
        tf.margin_right = int(Inches(0.06))
        tf.margin_top = int(Inches(0.08))
        tf.margin_bottom = int(Inches(0.05))

        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = label
        run.font.size = Pt(12)
        run.font.color.rgb = theme.TEXT_COLOR
        run.font.bold = True

        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = desc
        run2.font.size = Pt(10)
        run2.font.color.rgb = sec_color

        if i < n - 1:
            arr_x1 = int(x + box_w)
            arr_x2 = int(x + box_w + gap)
            arr_y = int(y + box_h // 2)
            arr = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT, arr_x1, arr_y, arr_x2, arr_y)
            arr.line.color.rgb = sec_color
            arr.line.width = Pt(2.0)
            _add_arrowhead(arr)

    if "note" in data:
        add_note(slide, data["note"])
    if data.get("_slide_num"):
        add_footer_bar(slide, data["_slide_num"], section)


def build_stack_diagram(slide, data):
    set_slide_background(slide, theme.BG_COLOR)
    section = data.get("section", 0)
    add_title_bar(slide, data["title"], section)

    sec_color = theme.SECTION_COLORS.get(section, theme.PRIMARY_COLOR) if section > 0 else theme.PRIMARY_COLOR
    layers = data["layers"]
    n = len(layers)
    box_h = int(Inches(0.42))   # 縮小：避免超出畫面
    gap = int(Inches(0.04))
    total_h = n * (box_h + gap) - gap
    # 垂直置中在 content area
    start_y = int(theme.CONTENT_TOP + (theme.CONTENT_HEIGHT - total_h) // 2)

    color_map = {
        "#2d5a27": RGBColor(0x2d, 0x5a, 0x27),
        "#1a3a6b": RGBColor(0x1a, 0x3a, 0x6b),
        "#4a2080": RGBColor(0x4a, 0x20, 0x80),
        "#7a1a1a": RGBColor(0x7a, 0x1a, 0x1a),
        "#3a3a00": RGBColor(0x3a, 0x3a, 0x00),
    }

    box_left = int(theme.CONTENT_LEFT + Inches(0.3))
    box_w = int(Inches(7.2))

    for i, layer_info in enumerate(layers):
        label, sublabel, color_hex = layer_info
        y = int(start_y + i * (box_h + gap))
        fill_color = color_map.get(color_hex, RGBColor(0x1a, 0x3a, 0x6b))

        shape = slide.shapes.add_shape(1, box_left, y, box_w, box_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.color.rgb = sec_color
        shape.line.width = Pt(1.5)

        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = label
        run.font.size = Pt(14)
        run.font.color.rgb = theme.TEXT_COLOR
        run.font.bold = True

        if sublabel:
            sub_left = int(box_left + box_w + Inches(0.1))
            sub_w = int(theme.CONTENT_RIGHT - sub_left)
            if sub_w > 0:
                add_textbox(slide, sublabel, sub_left, y, sub_w, box_h,
                            font_size=Pt(10), color=theme.SUBTEXT_COLOR)

        # 層間向下箭頭
        if i < n - 1:
            arr_x = int(box_left + box_w // 2)
            arr_y1 = int(y + box_h)
            arr_y2 = int(y + box_h + gap)
            arr = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT, arr_x, arr_y1, arr_x, arr_y2)
            arr.line.color.rgb = sec_color
            arr.line.width = Pt(1.5)

    if "note" in data:
        add_note(slide, data["note"])
    if data.get("_slide_num"):
        add_footer_bar(slide, data["_slide_num"], section)


def build_stack_diagram_annotated(slide, data):
    set_slide_background(slide, theme.BG_COLOR)
    section = data.get("section", 0)
    add_title_bar(slide, data["title"], section)

    layers = data["layers"]
    n = len(layers)
    box_h = int(Inches(0.44))   # 縮小：更安全
    gap = int(Inches(0.04))
    start_y = int(theme.CONTENT_TOP + Inches(0.05))

    status_fills = {
        "normal":  RGBColor(0x1a, 0x3a, 0x6b),
        "warning": RGBColor(0x6b, 0x4a, 0x00),
        "danger":  RGBColor(0x6b, 0x1a, 0x1a),
        "source":  RGBColor(0x1a, 0x5a, 0x27),
    }
    status_borders = {
        "normal":  theme.PRIMARY_COLOR,
        "warning": RGBColor(0xff, 0xcc, 0x44),
        "danger":  theme.ACCENT_COLOR,
        "source":  theme.SOURCE_ANNOTATION_COLOR,
    }

    box_left = int(theme.CONTENT_LEFT)
    box_w = int(Inches(5.2))
    ann_left = int(theme.CONTENT_LEFT + Inches(5.4))
    ann_w = int(theme.CONTENT_RIGHT - ann_left)

    for i, (label, annotation, status) in enumerate(layers):
        y = int(start_y + i * (box_h + gap))
        fill_color = status_fills.get(status, status_fills["normal"])
        border_color = status_borders.get(status, theme.PRIMARY_COLOR)
        border_w = Pt(2.5) if status in ("danger", "warning") else Pt(1.5)

        shape = slide.shapes.add_shape(1, box_left, y, box_w, box_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.color.rgb = border_color
        shape.line.width = border_w

        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = label
        run.font.size = Pt(14)
        run.font.color.rgb = theme.TEXT_COLOR
        run.font.bold = True

        ann_color = theme.ACCENT_COLOR if status in ("danger", "warning") else theme.SUBTEXT_COLOR
        if status == "source":
            ann_color = theme.SOURCE_ANNOTATION_COLOR
        if ann_w > 0:
            add_textbox(slide, annotation,
                        ann_left, y + int(Inches(0.05)), ann_w, box_h,
                        font_size=Pt(13), color=ann_color)

    if "note" in data:
        add_note(slide, data["note"])
    if data.get("_slide_num"):
        add_footer_bar(slide, data["_slide_num"], section)


def build_references(slide, data):
    set_slide_background(slide, theme.BG_COLOR)
    section = data.get("section", 0)
    add_title_bar(slide, data.get("qanda", "Q&A") + " / 參考資料", section)
    add_content_panel(slide)

    txBox = slide.shapes.add_textbox(
        int(theme.CONTENT_LEFT + Inches(0.2)),
        int(theme.CONTENT_TOP + Inches(0.15)),
        int(theme.CONTENT_WIDTH - Inches(0.4)),
        int(theme.CONTENT_HEIGHT - Inches(0.3))
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for ref in data["refs"]:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = ref
        run.font.size = Pt(13)
        run.font.color.rgb = theme.SUBTEXT_COLOR
        run.font.name = theme.FONT_BODY
        p.space_before = Pt(4)

    if data.get("_slide_num"):
        add_footer_bar(slide, data["_slide_num"], section)


BUILDERS = {
    "cover":                    build_cover,
    "section_break":            build_section_break,
    "bullets":                  build_bullets,
    "two_col":                  build_two_col,
    "table":                    build_table,
    "flow":                     build_flow,
    "stack_diagram":            build_stack_diagram,
    "stack_diagram_annotated":  build_stack_diagram_annotated,
    "references":               build_references,
}
