# src/create_ppt.py
import sys
import os
sys.path.insert(0, os.path.dirname(__file__))

from pptx import Presentation
from pptx.util import Inches
import theme
import content
import builders

OUTPUT_PATH = os.path.join(
    os.path.dirname(__file__), "..", "output", "network-card-csi-v5.3.pptx"
)

IMAGES_DIR = os.path.join(os.path.dirname(__file__), "..", "output", "images")

# {slide_index: (filename, left_in, top_in, width_in, height_in)}
# 16:9 比例（1.778:1）
# 注意：index 2 新增「為什麼需要 CSI？」投影片，後續 index 均 +1
SLIDE_IMAGES = {
    4:  ("slide03_rf_signal.png",    7.5,  1.10, 2.00, 1.13),  # 802.11 RF：flow 右上方空白區
    5:  ("slide04_csi_spectrum.png", 7.75, 3.30, 1.80, 1.01),  # PHY/MAC：右欄底部
    9:  ("slide08_nic_hardware.png", 7.10, 3.20, 2.30, 1.29),  # OEM.sys：右欄底部
    18: ("slide15_linux_open.png",   7.30, 3.20, 2.25, 1.27),  # Linux：右欄底部（+1 因新增破解OEM.sys slide）
}


def create_presentation():
    prs = Presentation()
    prs.slide_width = theme.SLIDE_WIDTH
    prs.slide_height = theme.SLIDE_HEIGHT

    blank_layout = prs.slide_layouts[6]  # completely blank layout

    for idx, slide_data in enumerate(content.SLIDES):
        enriched = {**slide_data, "_slide_num": idx + 1}
        slide = prs.slides.add_slide(blank_layout)
        slide_type = enriched["type"]
        builder_fn = builders.BUILDERS.get(slide_type)
        if builder_fn is None:
            print(f"WARNING: unknown slide type '{slide_type}', skipping")
            continue
        builder_fn(slide, enriched)
        if enriched.get("speaker_notes"):
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = enriched["speaker_notes"]

        # 插入 AI 生成圖片（若有設定且檔案存在）
        if idx in SLIDE_IMAGES:
            fname, l, t, w, h = SLIDE_IMAGES[idx]
            img_path = os.path.join(IMAGES_DIR, fname)
            if os.path.exists(img_path):
                builders.add_slide_image(
                    slide, img_path,
                    Inches(l), Inches(t), Inches(w), Inches(h)
                )
            else:
                print(f"WARNING: image not found: {img_path}")

    os.makedirs(os.path.dirname(OUTPUT_PATH), exist_ok=True)
    prs.save(OUTPUT_PATH)
    print(f"Saved: {OUTPUT_PATH}")
    print(f"Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    create_presentation()
